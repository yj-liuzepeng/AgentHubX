"""
文档生成器模块
将翻译后的内容重新生成各种格式的文档
"""

import os
import tempfile
from abc import ABC, abstractmethod
from typing import Dict, List, Optional, Any
from pathlib import Path
from loguru import logger

class BaseGenerator(ABC):
    """基础生成器类"""
    
    @abstractmethod
    def generate(self, content: Dict, output_path: str, original_path: str = None, preserve_formatting: bool = True) -> bool:
        """生成文档"""
        pass
    
    @abstractmethod
    def supports_format(self, format_type: str) -> bool:
        """检查是否支持该格式"""
        pass

class PDFGenerator(BaseGenerator):
    """PDF文档生成器"""
    
    def __init__(self):
        try:
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont
            
            self.letter = letter
            self.A4 = A4
            self.SimpleDocTemplate = SimpleDocTemplate
            self.Paragraph = Paragraph
            self.Spacer = Spacer
            self.PageBreak = PageBreak
            self.getSampleStyleSheet = getSampleStyleSheet
            self.ParagraphStyle = ParagraphStyle
            self.inch = inch
            self.pdfmetrics = pdfmetrics
            self.TTFont = TTFont
            
            # 注册中文字体
            try:
                self._register_chinese_fonts()
            except Exception as e:
                logger.warning(f"注册中文字体失败: {str(e)}")
                
        except ImportError:
            logger.error("reportlab库未安装，无法生成PDF文件")
            raise
    
    def supports_format(self, format_type: str) -> bool:
        return format_type.lower() == 'pdf'
    
    def _register_chinese_fonts(self):
        """注册中文字体"""
        # 尝试注册常用的中文字体
        chinese_fonts = [
            ('SimSun', 'simsun.ttc'),
            ('Microsoft YaHei', 'msyh.ttc'),
            ('SimHei', 'simhei.ttf')
        ]
        
        for font_name, font_file in chinese_fonts:
            try:
                if os.path.exists(f"C:/Windows/Fonts/{font_file}"):
                    self.pdfmetrics.registerFont(self.TTFont(font_name, f"C:/Windows/Fonts/{font_file}"))
                    logger.info(f"注册字体成功: {font_name}")
                    break
            except Exception as e:
                logger.warning(f"注册字体 {font_name} 失败: {str(e)}")
    
    def generate(self, content: Dict, output_path: str, original_path: str = None, preserve_formatting: bool = True) -> bool:
        """生成PDF文档"""
        try:
            # 创建PDF文档
            doc = self.SimpleDocTemplate(output_path, pagesize=self.A4)
            story = []
            
            # 获取样式
            styles = self.getSampleStyleSheet()
            
            # 创建自定义样式
            chinese_style = self.ParagraphStyle(
                'ChineseStyle',
                parent=styles['Normal'],
                fontName='SimSun',
                fontSize=12,
                spaceAfter=12,
                leading=14
            )
            
            title_style = self.ParagraphStyle(
                'TitleStyle',
                parent=styles['Title'],
                fontName='SimSun',
                fontSize=18,
                spaceAfter=20,
                alignment=1  # 居中
            )
            
            # 添加标题
            if 'metadata' in content and content['metadata'].get('title'):
                title = content['metadata']['title']
                story.append(self.Paragraph(title, title_style))
                story.append(self.Spacer(1, 12))
            
            # 添加翻译提示
            if 'translated_language' in content:
                notice_style = self.ParagraphStyle(
                    'NoticeStyle',
                    parent=styles['Normal'],
                    fontName='SimSun',
                    fontSize=10,
                    textColor='gray',
                    alignment=1  # 居中
                )
                notice = f"本文档已由系统自动翻译为 {content['translated_language'].upper()}"
                story.append(self.Paragraph(notice, notice_style))
                story.append(self.Spacer(1, 20))
            
            # 添加主要内容
            if preserve_formatting and 'pages' in content:
                # 保持页面格式
                for page in content['pages']:
                    if 'text' in page and page['text']:
                        paragraphs = page['text'].split('\n\n')
                        for para in paragraphs:
                            if para.strip():
                                story.append(self.Paragraph(para.strip(), chinese_style))
                                story.append(self.Spacer(1, 6))
                        story.append(self.PageBreak())
            else:
                # 简单文本格式
                text_content = content.get('text', '')
                if text_content:
                    paragraphs = text_content.split('\n\n')
                    for para in paragraphs:
                        if para.strip():
                            story.append(self.Paragraph(para.strip(), chinese_style))
                            story.append(self.Spacer(1, 6))
            
            # 生成PDF
            doc.build(story)
            logger.info(f"PDF文档生成成功: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"PDF文档生成失败: {str(e)}")
            return False

class DOCXGenerator(BaseGenerator):
    """DOCX文档生成器"""
    
    def __init__(self):
        try:
            from docx import Document
            from docx.shared import Pt, Inches, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.oxml.shared import OxmlElement, qn
            
            self.Document = Document
            self.Pt = Pt
            self.Inches = Inches
            self.RGBColor = RGBColor
            self.WD_ALIGN_PARAGRAPH = WD_ALIGN_PARAGRAPH
            self.OxmlElement = OxmlElement
            self.qn = qn
            
        except ImportError:
            logger.error("python-docx库未安装，无法生成DOCX文件")
            raise
    
    def supports_format(self, format_type: str) -> bool:
        return format_type.lower() == 'docx'
    
    def generate(self, content: Dict, output_path: str, original_path: str = None, preserve_formatting: bool = True) -> bool:
        """生成DOCX文档"""
        try:
            doc = self.Document()
            
            # 添加标题
            if 'metadata' in content and content['metadata'].get('title'):
                title = doc.add_heading(content['metadata']['title'], 0)
                title.alignment = self.WD_ALIGN_PARAGRAPH.CENTER
            
            # 添加翻译提示
            if 'translated_language' in content:
                notice_para = doc.add_paragraph()
                notice_run = notice_para.add_run(f"本文档已由系统自动翻译为 {content['translated_language'].upper()}")
                notice_run.font.color.rgb = self.RGBColor(128, 128, 128)
                notice_run.font.size = self.Pt(10)
                notice_para.alignment = self.WD_ALIGN_PARAGRAPH.CENTER
                doc.add_paragraph()  # 空行
            
            # 添加主要内容
            if preserve_formatting and 'paragraphs' in content:
                # 保持段落格式
                for para_info in content['paragraphs']:
                    if 'text' in para_info and para_info['text']:
                        para = doc.add_paragraph()
                        
                        # 设置段落样式
                        if 'style' in para_info:
                            if 'Heading' in para_info['style']:
                                level = 1 if '1' in para_info['style'] else 2
                                para = doc.add_heading(para_info['text'], level)
                            else:
                                para = doc.add_paragraph(para_info['text'])
                        else:
                            para = doc.add_paragraph(para_info['text'])
                        
                        # 设置对齐方式
                        if 'alignment' in para_info:
                            if para_info['alignment'] == 'CENTER':
                                para.alignment = self.WD_ALIGN_PARAGRAPH.CENTER
                            elif para_info['alignment'] == 'RIGHT':
                                para.alignment = self.WD_ALIGN_PARAGRAPH.RIGHT
                            else:
                                para.alignment = self.WD_ALIGN_PARAGRAPH.LEFT
                        
                        # 设置运行属性
                        if 'runs' in para_info and para_info['runs']:
                            para.clear()
                            for run_info in para_info['runs']:
                                if 'text' in run_info and run_info['text']:
                                    run = para.add_run(run_info['text'])
                                    
                                    # 设置字体样式
                                    if run_info.get('bold'):
                                        run.font.bold = True
                                    if run_info.get('italic'):
                                        run.font.italic = True
                                    if run_info.get('underline'):
                                        run.font.underline = True
                                    
                                    # 设置字体和大小
                                    if run_info.get('font_size'):
                                        run.font.size = self.Pt(run_info['font_size'])
            
            # 添加表格
            if 'tables' in content and content['tables']:
                for table_data in content['tables']:
                    if table_data:
                        rows = len(table_data)
                        cols = len(table_data[0]) if table_data else 0
                        
                        if rows > 0 and cols > 0:
                            table = doc.add_table(rows=rows, cols=cols)
                            table.style = 'Table Grid'
                            
                            for i, row_data in enumerate(table_data):
                                for j, cell_text in enumerate(row_data):
                                    cell = table.cell(i, j)
                                    cell.text = str(cell_text)
            
            # 保存文档
            doc.save(output_path)
            logger.info(f"DOCX文档生成成功: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"DOCX文档生成失败: {str(e)}")
            return False

class TXTGenerator(BaseGenerator):
    """TXT文档生成器"""
    
    def supports_format(self, format_type: str) -> bool:
        return format_type.lower() == 'txt'
    
    def generate(self, content: Dict, output_path: str, original_path: str = None, preserve_formatting: bool = True) -> bool:
        """生成TXT文档"""
        try:
            with open(output_path, 'w', encoding='utf-8') as f:
                # 添加标题
                if 'metadata' in content and content['metadata'].get('title'):
                    f.write(f"标题: {content['metadata']['title']}\n")
                    f.write("=" * 50 + "\n\n")
                
                # 添加翻译提示
                if 'translated_language' in content:
                    f.write(f"[本文档已由系统自动翻译为 {content['translated_language'].upper()}]\n\n")
                
                # 添加主要内容
                text_content = content.get('text', '')
                if text_content:
                    f.write(text_content)
                
                # 添加元数据信息
                if 'metadata' in content:
                    f.write("\n\n" + "=" * 50 + "\n")
                    f.write("文档信息:\n")
                    for key, value in content['metadata'].items():
                        if key != 'title' and value:
                            f.write(f"{key}: {value}\n")
            
            logger.info(f"TXT文档生成成功: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"TXT文档生成失败: {str(e)}")
            return False

class PPTGenerator(BaseGenerator):
    """PPT/PPTX文档生成器"""
    
    def __init__(self):
        try:
            from pptx import Presentation
            from pptx.util import Pt, Inches
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
            
            self.Presentation = Presentation
            self.Pt = Pt
            self.Inches = Inches
            self.PP_ALIGN = PP_ALIGN
            self.RGBColor = RGBColor
            
        except ImportError:
            logger.error("python-pptx库未安装，无法生成PPT文件")
            raise
    
    def supports_format(self, format_type: str) -> bool:
        return format_type.lower() in ['.ppt', '.pptx']
    
    def generate(self, content: Dict, output_path: str, original_path: str = None, preserve_formatting: bool = True) -> bool:
        """生成PPT文档"""
        try:
            prs = self.Presentation()
            
            # 添加标题页
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            # 设置标题
            if 'metadata' in content and content['metadata'].get('title'):
                title.text = content['metadata']['title']
            else:
                title.text = "翻译文档"
            
            # 设置副标题
            if 'translated_language' in content:
                subtitle.text = f"翻译为 {content['translated_language'].upper()}"
            else:
                subtitle.text = "自动翻译"
            
            # 添加内容幻灯片
            if 'slides' in content and content['slides']:
                for slide_info in content['slides']:
                    if 'text' in slide_info and slide_info['text']:
                        # 创建内容幻灯片
                        content_slide_layout = prs.slide_layouts[1]
                        slide = prs.slides.add_slide(content_slide_layout)
                        
                        # 设置标题
                        title_shape = slide.shapes.title
                        if 'title' in slide_info and slide_info['title']:
                            title_shape.text = slide_info['title']
                        else:
                            title_shape.text = "内容"
                        
                        # 设置内容
                        content_shape = slide.placeholders[1]
                        if 'content' in slide_info and slide_info['content']:
                            content_text = '\n'.join(slide_info['content'])
                            content_shape.text = content_text
                        
                        # 添加备注
                        if 'notes' in slide_info and slide_info['notes']:
                            notes_slide = slide.notes_slide
                            notes_slide.notes_text_frame.text = slide_info['notes']
            else:
                # 如果没有幻灯片结构，直接添加文本内容
                if 'text' in content and content['text']:
                    text_content = content['text']
                    paragraphs = text_content.split('\n\n')
                    
                    for para in paragraphs:
                        if para.strip():
                            content_slide_layout = prs.slide_layouts[1]
                            slide = prs.slides.add_slide(content_slide_layout)
                            
                            title_shape = slide.shapes.title
                            title_shape.text = "内容"
                            
                            content_shape = slide.placeholders[1]
                            content_shape.text = para.strip()
            
            # 保存演示文稿
            prs.save(output_path)
            logger.info(f"PPT文档生成成功: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"PPT文档生成失败: {str(e)}")
            return False

class DocumentGenerator:
    """文档生成器，支持多种格式"""
    
    def __init__(self):
        self.generators = {
            'pdf': PDFGenerator(),
            'docx': DOCXGenerator(),
            'txt': TXTGenerator(),
            'ppt': PPTGenerator(),
            'pptx': PPTGenerator()
        }
    
    def generate(self, content: Dict, output_path: str, original_path: str = None, preserve_formatting: bool = True) -> bool:
        """
        生成文档
        
        Args:
            content: 文档内容
            output_path: 输出文件路径
            original_path: 原始文件路径（用于格式参考）
            preserve_formatting: 是否保持格式
        
        Returns:
            是否成功生成
        """
        try:
            # 根据输出文件扩展名选择生成器
            output_ext = Path(output_path).suffix.lower()
            
            if output_ext == '.pdf':
                generator = self.generators['pdf']
            elif output_ext == '.docx':
                generator = self.generators['docx']
            elif output_ext == '.txt':
                generator = self.generators['txt']
            elif output_ext in ['.ppt', '.pptx']:
                generator = self.generators['ppt']
            else:
                logger.error(f"不支持的输出格式: {output_ext}")
                return False
            
            # 生成文档
            return generator.generate(content, output_path, original_path, preserve_formatting)
            
        except Exception as e:
            logger.error(f"文档生成失败: {str(e)}")
            return False
    
    def get_supported_formats(self) -> List[str]:
        """获取支持的输出格式"""
        return ['.pdf', '.docx', '.txt', '.ppt', '.pptx']