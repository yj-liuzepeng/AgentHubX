"""
文档解析器模块
支持PDF、DOCX、DOC、TXT、PPT格式文件的文本提取
"""

import os
import re
from abc import ABC, abstractmethod
from typing import Dict, List, Optional
from pathlib import Path
from loguru import logger

class BaseParser(ABC):
    """基础解析器类"""
    
    @abstractmethod
    def parse(self, file_path: str) -> Optional[Dict]:
        """解析文件并返回结构化内容"""
        pass
    
    @abstractmethod
    def supports_format(self, file_extension: str) -> bool:
        """检查是否支持该文件格式"""
        pass

class PDFParser(BaseParser):
    """PDF文件解析器"""
    
    def __init__(self):
        try:
            import fitz  # PyMuPDF
            self.fitz = fitz
        except ImportError:
            logger.error("PyMuPDF库未安装，无法解析PDF文件")
            raise
    
    def supports_format(self, file_extension: str) -> bool:
        return file_extension.lower() == '.pdf'
    
    def parse(self, file_path: str) -> Optional[Dict]:
        """解析PDF文件"""
        try:
            doc = self.fitz.open(file_path)
            content = {
                'text': '',
                'pages': [],
                'metadata': {},
                'format': 'pdf'
            }
            
            # 提取元数据
            metadata = doc.metadata
            content['metadata'] = {
                'title': metadata.get('title', ''),
                'author': metadata.get('author', ''),
                'subject': metadata.get('subject', ''),
                'creator': metadata.get('creator', ''),
                'creation_date': metadata.get('creationDate', ''),
                'total_pages': len(doc)
            }
            
            # 提取每页内容
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                page_text = page.get_text()
                
                # 提取页面信息
                page_info = {
                    'page_number': page_num + 1,
                    'text': page_text,
                    'rect': list(page.rect),
                    'blocks': self._extract_blocks(page)
                }
                
                content['pages'].append(page_info)
                content['text'] += page_text + '\n\n'
            
            doc.close()
            
            # 清理文本
            content['text'] = self._clean_text(content['text'])
            
            return content
            
        except Exception as e:
            logger.error(f"PDF解析失败: {str(e)}")
            return None
    
    def _extract_blocks(self, page) -> List[Dict]:
        """提取页面中的文本块"""
        try:
            blocks = page.get_text("dict")["blocks"]
            block_info = []
            
            for block in blocks:
                if "lines" in block:  # 文本块
                    block_text = ""
                    for line in block["lines"]:
                        for span in line["spans"]:
                            block_text += span["text"] + " "
                    
                    block_info.append({
                        'type': 'text',
                        'text': block_text.strip(),
                        'bbox': block["bbox"],
                        'font_size': block["lines"][0]["spans"][0]["size"] if block["lines"] else 12
                    })
            
            return block_info
            
        except Exception as e:
            logger.warning(f"提取文本块失败: {str(e)}")
            return []
    
    def _clean_text(self, text: str) -> str:
        """清理文本内容"""
        # 移除多余的空行和空格
        text = re.sub(r'\n{3,}', '\n\n', text)
        text = re.sub(r'[ \t]+', ' ', text)
        return text.strip()

class DOCXParser(BaseParser):
    """DOCX文件解析器"""
    
    def __init__(self):
        try:
            from docx import Document
            self.Document = Document
        except ImportError:
            logger.error("python-docx库未安装，无法解析DOCX文件")
            raise
    
    def supports_format(self, file_extension: str) -> bool:
        return file_extension.lower() == '.docx'
    
    def parse(self, file_path: str) -> Optional[Dict]:
        """解析DOCX文件"""
        try:
            doc = self.Document(file_path)
            content = {
                'text': '',
                'paragraphs': [],
                'tables': [],
                'metadata': {},
                'format': 'docx'
            }
            
            # 提取段落
            for para in doc.paragraphs:
                para_text = para.text.strip()
                if para_text:
                    para_info = {
                        'text': para_text,
                        'style': para.style.name if para.style else 'Normal',
                        'alignment': str(para.alignment) if para.alignment else 'LEFT',
                        'runs': self._extract_runs(para)
                    }
                    content['paragraphs'].append(para_info)
                    content['text'] += para_text + '\n\n'
            
            # 提取表格
            for table in doc.tables:
                table_data = self._extract_table(table)
                if table_data:
                    content['tables'].append(table_data)
            
            # 提取文档属性
            core_props = doc.core_properties
            content['metadata'] = {
                'title': core_props.title or '',
                'author': core_props.author or '',
                'subject': core_props.subject or '',
                'created': str(core_props.created) if core_props.created else '',
                'total_paragraphs': len(content['paragraphs']),
                'total_tables': len(content['tables'])
            }
            
            # 清理文本
            content['text'] = self._clean_text(content['text'])
            
            return content
            
        except Exception as e:
            logger.error(f"DOCX解析失败: {str(e)}")
            return None
    
    def _extract_runs(self, paragraph) -> List[Dict]:
        """提取段落中的文本运行"""
        runs = []
        for run in paragraph.runs:
            if run.text.strip():
                runs.append({
                    'text': run.text,
                    'bold': run.bold,
                    'italic': run.italic,
                    'underline': run.underline,
                    'font_name': run.font.name if run.font.name else 'Arial',
                    'font_size': run.font.size.pt if run.font.size else 12
                })
        return runs
    
    def _extract_table(self, table) -> List[List[str]]:
        """提取表格数据"""
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                row_data.append(cell_text)
            table_data.append(row_data)
        return table_data
    
    def _clean_text(self, text: str) -> str:
        """清理文本内容"""
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text.strip()

class DOCParser(BaseParser):
    """DOC文件解析器"""
    
    def __init__(self):
        try:
            import win32com.client
            self.win32com = win32com.client
        except ImportError:
            logger.warning("win32com库未安装，将使用antiword作为DOC文件解析备选方案")
            self.win32com = None
    
    def supports_format(self, file_extension: str) -> bool:
        return file_extension.lower() == '.doc'
    
    def parse(self, file_path: str) -> Optional[Dict]:
        """解析DOC文件"""
        try:
            if self.win32com:
                return self._parse_with_win32com(file_path)
            else:
                return self._parse_with_antiword(file_path)
                
        except Exception as e:
            logger.error(f"DOC解析失败: {str(e)}")
            return None
    
    def _parse_with_win32com(self, file_path: str) -> Optional[Dict]:
        """使用win32com解析DOC文件"""
        try:
            word = self.win32com.Dispatch("Word.Application")
            word.Visible = False
            
            doc = word.Documents.Open(file_path)
            content = {
                'text': doc.Content.Text,
                'paragraphs': [],
                'metadata': {},
                'format': 'doc'
            }
            
            # 提取段落
            for para in doc.Paragraphs:
                para_text = para.Range.Text.strip()
                if para_text:
                    content['paragraphs'].append({
                        'text': para_text,
                        'style': para.Style.NameLocal if para.Style else 'Normal'
                    })
            
            # 提取文档属性
            content['metadata'] = {
                'title': doc.BuiltInDocumentProperties("Title") or '',
                'author': doc.BuiltInDocumentProperties("Author") or '',
                'total_paragraphs': len(content['paragraphs'])
            }
            
            doc.Close()
            word.Quit()
            
            content['text'] = self._clean_text(content['text'])
            return content
            
        except Exception as e:
            logger.error(f"win32com解析DOC失败: {str(e)}")
            return None
    
    def _parse_with_antiword(self, file_path: str) -> Optional[Dict]:
        """使用antiword解析DOC文件"""
        try:
            import subprocess
            
            # 使用antiword提取文本
            result = subprocess.run(
                ['antiword', file_path],
                capture_output=True,
                text=True,
                encoding='utf-8'
            )
            
            if result.returncode == 0:
                text = result.stdout
                content = {
                    'text': text,
                    'paragraphs': [{'text': para.strip()} for para in text.split('\n\n') if para.strip()],
                    'metadata': {'format': 'doc', 'parser': 'antiword'},
                    'format': 'doc'
                }
                
                content['text'] = self._clean_text(content['text'])
                return content
            else:
                logger.error(f"antiword解析失败: {result.stderr}")
                return None
                
        except Exception as e:
            logger.error(f"antiword解析DOC失败: {str(e)}")
            return None
    
    def _clean_text(self, text: str) -> str:
        """清理文本内容"""
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text.strip()

class TXTParser(BaseParser):
    """TXT文件解析器"""
    
    def supports_format(self, file_extension: str) -> bool:
        return file_extension.lower() == '.txt'
    
    def parse(self, file_path: str) -> Optional[Dict]:
        """解析TXT文件"""
        try:
            # 检测文件编码
            encoding = self._detect_encoding(file_path)
            
            with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
                text_content = f.read()
            
            content = {
                'text': text_content,
                'paragraphs': [{'text': para.strip()} for para in text_content.split('\n\n') if para.strip()],
                'metadata': {
                    'encoding': encoding,
                    'line_count': len(text_content.split('\n')),
                    'char_count': len(text_content)
                },
                'format': 'txt'
            }
            
            content['text'] = self._clean_text(content['text'])
            return content
            
        except Exception as e:
            logger.error(f"TXT解析失败: {str(e)}")
            return None
    
    def _detect_encoding(self, file_path: str) -> str:
        """检测文件编码"""
        try:
            import chardet
            
            with open(file_path, 'rb') as f:
                raw_data = f.read(10000)  # 读取前10KB
            
            result = chardet.detect(raw_data)
            encoding = result['encoding'] or 'utf-8'
            
            # 验证编码
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    f.read(1000)
                return encoding
            except:
                return 'utf-8'
                
        except ImportError:
            logger.warning("chardet库未安装，使用默认UTF-8编码")
            return 'utf-8'
        except Exception as e:
            logger.error(f"编码检测失败: {str(e)}")
            return 'utf-8'
    
    def _clean_text(self, text: str) -> str:
        """清理文本内容"""
        # 移除多余的空行
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text.strip()

class PPTParser(BaseParser):
    """PPT/PPTX文件解析器"""
    
    def __init__(self):
        try:
            from pptx import Presentation
            self.Presentation = Presentation
        except ImportError:
            logger.error("python-pptx库未安装，无法解析PPT文件")
            raise
    
    def supports_format(self, file_extension: str) -> bool:
        return file_extension.lower() in ['.ppt', '.pptx']
    
    def parse(self, file_path: str) -> Optional[Dict]:
        """解析PPT文件"""
        try:
            prs = self.Presentation(file_path)
            content = {
                'text': '',
                'slides': [],
                'metadata': {},
                'format': 'ppt'
            }
            
            # 提取每张幻灯片
            for slide_idx, slide in enumerate(prs.slides):
                slide_text = []
                slide_info = {
                    'slide_number': slide_idx + 1,
                    'title': '',
                    'content': [],
                    'notes': ''
                }
                
                # 提取标题
                title_shape = None
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        if shape == slide.shapes.title:
                            title_shape = shape
                            slide_info['title'] = shape.text_frame.text
                        else:
                            slide_text.append(shape.text_frame.text)
                
                # 提取内容
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame and shape != title_shape:
                        text = shape.text_frame.text.strip()
                        if text:
                            slide_info['content'].append(text)
                
                # 提取备注
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame:
                        slide_info['notes'] = notes_slide.notes_text_frame.text
                
                content['slides'].append(slide_info)
                
                # 合并所有文本
                all_slide_text = slide_info['title'] + '\n' + '\n'.join(slide_info['content'])
                if all_slide_text.strip():
                    content['text'] += all_slide_text + '\n\n'
            
            # 提取文档属性
            core_props = prs.core_properties
            content['metadata'] = {
                'title': core_props.title or '',
                'author': core_props.author or '',
                'subject': core_props.subject or '',
                'total_slides': len(prs.slides)
            }
            
            content['text'] = self._clean_text(content['text'])
            return content
            
        except Exception as e:
            logger.error(f"PPT解析失败: {str(e)}")
            return None
    
    def _clean_text(self, text: str) -> str:
        """清理文本内容"""
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text.strip()

# 解析器工厂
class ParserFactory:
    """解析器工厂类"""
    
    def __init__(self):
        self.parsers = [
            PDFParser(),
            DOCXParser(),
            DOCParser(),
            TXTParser(),
            PPTParser()
        ]
    
    def get_parser(self, file_extension: str) -> Optional[BaseParser]:
        """根据文件扩展名获取合适的解析器"""
        for parser in self.parsers:
            if parser.supports_format(file_extension):
                return parser
        return None
    
    def get_supported_formats(self) -> List[str]:
        """获取支持的文件格式列表"""
        formats = []
        for parser in self.parsers:
            # 这里简化处理，实际需要根据每个解析器获取具体格式
            if isinstance(parser, PDFParser):
                formats.append('.pdf')
            elif isinstance(parser, DOCXParser):
                formats.append('.docx')
            elif isinstance(parser, DOCParser):
                formats.append('.doc')
            elif isinstance(parser, TXTParser):
                formats.append('.txt')
            elif isinstance(parser, PPTParser):
                formats.extend(['.ppt', '.pptx'])
        return formats